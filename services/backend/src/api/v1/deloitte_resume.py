import io
import os
import re
import tempfile
from pathlib import Path
from typing import Optional
from uuid import UUID

import httpx
import structlog
from fastapi import APIRouter, Depends, HTTPException
from pydantic import BaseModel
from sqlalchemy import select, and_
from sqlalchemy.ext.asyncio import AsyncSession

from src.db.models import CandidateCV, Candidate, ParsedResume
from src.db.session import get_db_session
from src.services.storage import upload_deloitte_pptx, delete_deloitte_pptx_from_storage

logger = structlog.get_logger()
router = APIRouter()


def _safe_str(val) -> str:
    if val is None:
        return ""
    s = str(val).strip()
    s = re.sub(r'[\x00-\x08\x0b\x0c\x0e-\x1f\x7f-\x9f]', '', s)
    s = s.replace('\x00', '')
    return s


async def _fetch_cv_text(file_url: str) -> str:
    try:
        async with httpx.AsyncClient(timeout=60.0) as client:
            resp = await client.get(file_url)
            resp.raise_for_status()
            content = resp.content

        file_url_lower = file_url.lower().split("?")[0]

        if file_url_lower.endswith(".pdf"):
            return _extract_text_from_pdf(content)
        elif file_url_lower.endswith(".docx"):
            return _extract_text_from_docx(content)
        elif file_url_lower.endswith(".pptx") or file_url_lower.endswith(".ppt"):
            return _extract_text_from_pptx(content)
        elif file_url_lower.endswith(".doc"):
            return _extract_text_from_docx(content)
        else:
            try:
                return _extract_text_from_pdf(content)
            except Exception:
                return content.decode("utf-8", errors="ignore")
    except Exception as e:
        logger.warning("cv_fetch_failed", error=str(e), url=file_url)
        return ""


def _extract_text_from_pdf(content: bytes) -> str:
    try:
        import pypdf
        reader = pypdf.PdfReader(io.BytesIO(content))
        parts = []
        for page in reader.pages:
            text = page.extract_text()
            if text:
                parts.append(text)
        return "\n\n".join(parts)
    except Exception:
        pass
    try:
        import pdfplumber
        with pdfplumber.open(io.BytesIO(content)) as pdf:
            parts = []
            for page in pdf.pages:
                text = page.extract_text()
                if text:
                    parts.append(text)
            return "\n\n".join(parts)
    except Exception:
        return content.decode("utf-8", errors="ignore")


def _extract_text_from_docx(content: bytes) -> str:
    try:
        from docx import Document
        doc = Document(io.BytesIO(content))
        parts = []
        for para in doc.paragraphs:
            text = para.text.strip()
            if text:
                parts.append(text)
        for table in doc.tables:
            for row in table.rows:
                row_text = " | ".join(cell.text.strip() for cell in row.cells if cell.text.strip())
                if row_text:
                    parts.append(row_text)
        return "\n\n".join(parts)
    except Exception:
        return content.decode("utf-8", errors="ignore")


def _extract_text_from_pptx(content: bytes) -> str:
    try:
        from pptx import Presentation
        prs = Presentation(io.BytesIO(content))
        parts = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    parts.append(shape.text.strip())
        return "\n\n".join(parts)
    except Exception:
        return content.decode("utf-8", errors="ignore")


def _parse_cv_text(cv_text: str, candidate: Candidate, parsed: Optional[ParsedResume]) -> dict:
    lines = [_safe_str(l) for l in cv_text.replace("\r\n", "\n").replace("\r", "\n").split("\n") if l.strip()]

    first = _safe_str(candidate.first_name)
    last = _safe_str(candidate.last_name)
    full_name = f"{first} {last}".strip()

    role = _safe_str(candidate.current_title)
    if not role:
        for line in lines[:5]:
            if line and line != full_name and "@" not in line and "+" not in line and len(line) > 3 and len(line) < 80 and "," not in line[:20]:
                role = line
                break

    location = _safe_str(candidate.location) if candidate.location else ""
    if not location:
        location_keywords = [
            "san ramon", "san francisco", "palo alto", "santa clara",
            "sunnyvale", "new york", "chicago", "seattle", "boston",
            "bangalore", "london", "remote", "austin", "dallas"
        ]
        for line in lines[:8]:
            if "@" in line or "+" in line or len(line) > 60:
                continue
            ll = line.lower()
            if any(kw in ll for kw in location_keywords):
                location = line.strip()
                break

    summary = ""
    if parsed and parsed.summary:
        summary = _safe_str(parsed.summary)

    if not summary:
        in_summary = False
        summary_lines_raw = []
        for line in lines:
            ll = line.lower()
            if ll.startswith("professional summary") or ll == "summary":
                in_summary = True
                continue
            if in_summary:
                if any(ll.startswith(s) for s in [
                    "core ", "competencies", "certifications",
                    "education", "professional experience"
                ]):
                    break
                if len(line) > 20:
                    summary_lines_raw.append(line)
        summary = " ".join(summary_lines_raw)

    summary_paras = []
    if summary:
        sentences = [s.strip() for s in summary.replace("\n", " ").split(".") if s.strip() and len(s.strip()) > 5]
        mid = max(1, len(sentences) // 2)
        para1 = ". ".join(sentences[:mid]).strip()
        if para1 and not para1.endswith("."):
            para1 += "."
        para2 = ". ".join(sentences[mid:]).strip()
        if para2 and not para2.endswith("."):
            para2 += "."
        if para1:
            summary_paras.append(_safe_str(para1[:350]))
        if para2:
            summary_paras.append(_safe_str(para2[:350]))

    if not summary_paras:
        summary_paras = [
            "A seasoned professional with extensive experience across multiple technical domains.",
            "Demonstrates a proven ability to deliver results and lead cross-functional teams."
        ]

    tech_skills_grouped = []
    business_skills = []

    TECH_ONLY_LABELS = [
        "security integration", "ci/cd", "container", "cloud platform",
        "observability", "storage"
    ]
    BUSINESS_ONLY_LABELS = ["leadership"]

    in_competencies = False
    for line in lines:
        ll = line.lower()

        if "competencies" in ll or "core devsecops" in ll:
            in_competencies = True
            continue

        if in_competencies:
            if any(ll.startswith(s) for s in [
                "certifications", "education", "professional experience"
            ]):
                in_competencies = False
                continue

            if ":" in line:
                parts = line.split(":", 1)
                label = parts[0].strip()
                detail = parts[1].strip() if len(parts) > 1 else ""
                label_lower = label.lower()

                if any(bl in label_lower for bl in BUSINESS_ONLY_LABELS):
                    items = [_safe_str(x.strip()) for x in detail.replace("|", ",").split(",") if x.strip()]
                    business_skills.extend(items[:5])
                elif any(tl in label_lower for tl in TECH_ONLY_LABELS):
                    if detail:
                        detail_items = [_safe_str(x.strip()) for x in detail.split(",") if x.strip()]
                        shortened = ", ".join(detail_items[:6])
                        tech_skills_grouped.append([_safe_str(label), _safe_str(shortened)])

    if not business_skills:
        skills_all = list(candidate.skills or [])
        BUSINESS_KW = [
            "agile", "scrum", "kanban", "leadership", "management",
            "capex", "c-level", "finops", "cross-functional", "program", "stakeholder"
        ]
        for s in skills_all:
            if any(kw in s.lower() for kw in BUSINESS_KW):
                business_skills.append(_safe_str(s))

    if not business_skills:
        business_skills = [
            "Cross-Functional Team Leadership",
            "Agile / Scrum / Kanban",
            "Stakeholder Management",
            "Program Delivery",
            "Strategic Planning",
        ]

    business_skills = [_safe_str(s) for s in business_skills[:5]]

    if not tech_skills_grouped:
        skills_all = list(candidate.skills or [])
        if skills_all:
            chunk_size = max(3, len(skills_all) // 5)
            chunks = [skills_all[i:i+chunk_size] for i in range(0, len(skills_all), chunk_size)]
            tech_skills_grouped = [
                [f"Technical Skills {i+1}", _safe_str(", ".join(chunk[:5]))]
                for i, chunk in enumerate(chunks[:6])
            ]

    languages = ["English"]
    if parsed and parsed.json_data and isinstance(parsed.json_data.get("languages"), list):
        languages = [_safe_str(l) for l in parsed.json_data["languages"]]

    cv_lower = cv_text.lower()
    INDUSTRY_MAP = [
        ("verizon", "Technology / Cloud Services"),
        ("satellite", "Defence / Aerospace"),
        ("iridium", "Defence / Aerospace"),
        ("aerospace", "Defence / Aerospace"),
        ("costco", "Retail"),
        ("retail", "Retail"),
        ("storage", "Storage / Infrastructure"),
        ("banking", "Financial Services"),
        ("finance", "Financial Services"),
        ("healthcare", "Healthcare"),
        ("manufacturing", "Manufacturing"),
        ("telecom", "Telecommunications"),
    ]

    industry_experience = []
    seen = set()
    for kw, label in INDUSTRY_MAP:
        if kw in cv_lower and label not in seen:
            industry_experience.append(label)
            seen.add(label)
        if len(industry_experience) >= 4:
            break

    if not industry_experience:
        industry_experience = ["Technology / Cloud Services"]

    education_lines = []
    in_edu = False
    for line in lines:
        ll = line.lower()
        if ll.startswith("education"):
            in_edu = True
            continue
        if in_edu:
            if any(ll.startswith(s) for s in [
                "professional experience", "experience", "certifications"
            ]):
                break
            if len(line) > 5:
                cleaned = _safe_str(line.split("\u2013")[0].strip())
                cleaned = cleaned.split(" - ")[0].strip() if " - " in cleaned else cleaned
                education_lines.append(cleaned[:50])
        if len(education_lines) >= 2:
            break

    cert_lines = []
    in_cert = False
    for line in lines:
        ll = line.lower()
        if ll.startswith("certifications"):
            in_cert = True
            continue
        if in_cert:
            if any(ll.startswith(s) for s in [
                "education", "professional experience", "experience", "core "
            ]):
                break
            if len(line) > 3 and "|" not in line:
                cleaned = _safe_str(line.replace("Certified Engineer \u2013 ", "").replace("Certified Engineer - ", "").strip())
                cleaned = cleaned.replace(" Course", "").strip()
                cert_lines.append(cleaned[:50])
        if len(cert_lines) >= 3:
            break

    clients = []
    relevant_exp_paras = []

    exp_raw = parsed.experience if parsed and parsed.experience else []
    if not exp_raw and candidate.json_data:
        exp_raw = candidate.json_data.get("experience", [])

    if exp_raw:
        for exp in exp_raw:
            if isinstance(exp, dict):
                company = _safe_str(exp.get("company", "")).split(",")[0].strip()
                if company and company not in clients:
                    clients.append(company)
                resp = exp.get("responsibilities", [])
                if isinstance(resp, list):
                    resp_text = ". ".join([_safe_str(r) for r in resp[:4] if _safe_str(r)])
                elif isinstance(resp, str):
                    resp_text = _safe_str(resp.strip())
                else:
                    resp_text = ""
                description = _safe_str(exp.get("description", ""))
                content = resp_text or description
                if content and len(relevant_exp_paras) < 3:
                    prefixes = ["Experience includes", "Skilled in", "Strong background in"]
                    para = _safe_str(f"{prefixes[len(relevant_exp_paras)]} {content[:350]}")
                    if not para.endswith("."):
                        last_period = para.rfind(".")
                        if last_period > 50:
                            para = para[:last_period + 1]
                        else:
                            para += "."
                    relevant_exp_paras.append(para)

    if not relevant_exp_paras:
        in_exp_section = False
        exp_resp_lines = []
        current_block = []

        EXP_HEADERS = [
            "professional experience", "work experience", "experience",
            "employment history", "career history"
        ]

        for line in lines:
            ll = line.lower()
            if any(ll.startswith(h) for h in EXP_HEADERS):
                in_exp_section = True
                continue
            if in_exp_section:
                date_pattern = re.search(
                    r'(jan|feb|mar|apr|may|jun|jul|aug|sep|oct|nov|dec)\s+\d{4}|'
                    r'\d{4}\s*[-\u2013]\s*\d{4}|'
                    r'\d{4}\s*[-\u2013]\s*present', ll
                )
                if date_pattern and len(line) < 120:
                    company_part = line.split(",")[0].strip()
                    if company_part and len(company_part) > 2 and company_part not in clients:
                        clients.append(company_part)
                    if current_block:
                        exp_resp_lines.append(" ".join(current_block))
                        current_block = []
                elif len(line) > 20:
                    current_block.append(line)

        if current_block:
            exp_resp_lines.append(" ".join(current_block))

        prefixes = ["Experience includes", "Skilled in", "Strong background in"]
        for i, block in enumerate(exp_resp_lines[:3]):
            para = _safe_str(f"{prefixes[i]} {block[:350]}")
            if not para.endswith("."):
                last_period = para.rfind(".")
                if last_period > 50:
                    para = para[:last_period + 1]
                else:
                    para += "."
            relevant_exp_paras.append(para)

    if not clients:
        in_exp_section = False
        for line in lines:
            ll = line.lower()
            if any(ll.startswith(h) for h in ["professional experience", "work experience", "experience"]):
                in_exp_section = True
                continue
            if in_exp_section:
                date_pattern = re.search(
                    r'(jan|feb|mar|apr|may|jun|jul|aug|sep|oct|nov|dec)\s+\d{4}|'
                    r'\d{4}\s*[-\u2013]\s*\d{4}', ll
                )
                if date_pattern and len(line) < 120:
                    company_part = line.split(",")[0].strip()
                    if company_part and len(company_part) > 2 and company_part not in clients:
                        clients.append(company_part)
            if len(clients) >= 8:
                break

    while len(relevant_exp_paras) < 3:
        relevant_exp_paras.append("")

    CLIENT_SHORT = {
        "VERIZON": "Verizon",
        "IRIDIUM SATELLITE COMM / CDW": "Iridium Satellite Comm / CDW",
        "COSTCO WHOLESALE": "Costco Wholesale",
        "ODYSSEY SPACE RESEARCH": "Odyssey Space Research",
        "AQUILA CLOUDS": "Aquila Clouds",
        "VERITAS TECHNOLOGIES LLC": "Veritas Technologies",
        "COHO DATA": "Coho Data",
        "HP HELION (OpenStack)": "HP Helion",
        "NETWORK APPLIANCE (NetApp)": "NetApp",
        "ATTUNE SYSTEMS": "Attune Systems",
    }

    clean_clients = []
    for c in clients[:8]:
        clean_clients.append(_safe_str(CLIENT_SHORT.get(c, c.title())))

    return {
        "name_large": full_name,
        "name_card": full_name,
        "role": f"Role: {role}" if role else "Role: Consultant",
        "location": f"Location: {location}" if location else "Location: Not specified",
        "summary_paras": summary_paras,
        "business_skills": business_skills,
        "tech_skills": tech_skills_grouped,
        "languages": languages,
        "industry_experience": industry_experience,
        "certifications": cert_lines,
        "education": education_lines,
        "relevant_exp_paras": relevant_exp_paras,
        "clients": ", ".join(clean_clients),
    }


def _estimate_text_height_emu(texts: list, size_pt: float, box_width_emu: int) -> int:
    EMU_PER_PT = 12700
    chars_per_line = max(1, int(box_width_emu / (size_pt * EMU_PER_PT * 0.55)))
    LINE_H = int(size_pt * EMU_PER_PT * 1.2)
    PARA_GAP = int(size_pt * EMU_PER_PT * 0.4)
    total = 0
    for text in texts:
        if not text or not text.strip():
            total += PARA_GAP
            continue
        n_lines = max(1, -(-len(text) // chars_per_line))
        total += n_lines * LINE_H + PARA_GAP
    return max(total, int(0.3 * 914400))


def _generate_pptx_bytes(candidate_data: dict, template_path: str) -> bytes:
    import shutil
    import lxml.etree as etree
    from pptx import Presentation

    if not Path(template_path).exists():
        raise FileNotFoundError(f"Template not found: {template_path}")

    with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as tmp:
        out_path = tmp.name

    try:
        shutil.copy(template_path, out_path)
        prs = Presentation(out_path)
        slide = prs.slides[0]

        for master_shape in prs.slide_master.shapes:
            if master_shape.name == "TextBox 2" and master_shape.shape_type == 17:
                master_shape._element.getparent().remove(master_shape._element)
                break

        SLIDE_HEIGHT = prs.slide_height

        PH26_TOP = 1548082
        PH26_LEFT = 6339600
        PH26_WIDTH = 5357100
        PH26_MIN_H = int(0.4 * 914400)
        PH26_MAX_H = 1127242

        PH16_TOP = 3081163
        PH16_LEFT = 925888
        PH16_WIDTH = 2304000
        PH16_MIN_H = int(0.5 * 914400)

        PH24_TOP = 3081163
        PH24_LEFT = 3348000
        PH24_WIDTH = 2304000
        PH24_MIN_H = int(0.5 * 914400)

        LEFT_BOTTOM_LIMIT = int(6.570 * 914400)

        EXP_TB_TOP = 3036417
        EXP_TB_LEFT = 6339600
        EXP_TB_WIDTH = 5357100
        EXP_TB_MAX_H = 2573514

        GROUP25_HEIGHT = 367041
        GROUP25_LEFT = 5913061
        GROUP25_WIDTH = 5890679

        CLIENTS_LEFT = 6339600
        CLIENTS_WIDTH = 5357100
        CLIENTS_HEIGHT = 324000

        SECTION_GAP = int(0.04 * 914400)

        EMU_PER_PT = 12700

        nsmap = "http://schemas.openxmlformats.org/drawingml/2006/main"

        def _apply_autofit(txBody):
            bodyPr = txBody.find(f'{{{nsmap}}}bodyPr')
            if bodyPr is None:
                return
            for child in list(bodyPr):
                tag = child.tag.split("}")[-1] if "}" in child.tag else child.tag
                if tag in ("noAutofit", "normAutofit", "spAutoFit"):
                    bodyPr.remove(child)
            etree.SubElement(bodyPr, f'{{{nsmap}}}normAutofit')

        def _estimate_h(texts: list, size_pt: float, box_width_emu: int) -> int:
            chars_per_line = max(1, int(box_width_emu / (size_pt * EMU_PER_PT * 0.52)))
            line_h = int(size_pt * EMU_PER_PT * 1.15)
            gap = int(size_pt * EMU_PER_PT * 0.5)
            total = 0
            for text in texts:
                if not text or not text.strip():
                    total += gap
                    continue
                n = max(1, -(-len(text) // chars_per_line))
                total += n * line_h + gap
            return max(total, int(0.3 * 914400))

        def _sanitize(text) -> str:
            if text is None:
                return ""
            s = str(text)
            s = re.sub(r'[\x00-\x08\x0b\x0c\x0e-\x1f\x7f-\x9f]', '', s)
            s = s.replace('\x00', '')
            return s

        def _truncate_to_fit(text: str, size_pt: float, box_width_emu: int, max_lines: int) -> str:
            EMU_PER_PT = 12700
            chars_per_line = max(1, int(box_width_emu / (size_pt * EMU_PER_PT * 0.52)))
            max_chars = chars_per_line * max_lines
            if len(text) <= max_chars:
                return text
            truncated = text[:max_chars - 3]
            last_space = truncated.rfind(' ')
            if last_space > max_chars * 0.8:
                truncated = truncated[:last_space]
            return truncated + "..."

        def _add_para(txBody, text: str, bold=None, size_pt=None, color_rgb=None, space_before=0, space_after=0):
            text = _sanitize(text)
            p_el = etree.SubElement(txBody, f'{{{nsmap}}}p')
            pPr = etree.SubElement(p_el, f'{{{nsmap}}}pPr')
            pPr.set("algn", "just")
            spcBef_el = etree.SubElement(pPr, f'{{{nsmap}}}spcBef')
            etree.SubElement(spcBef_el, f'{{{nsmap}}}spcPts').set("val", str(space_before * 100))
            spcAft_el = etree.SubElement(pPr, f'{{{nsmap}}}spcAft')
            etree.SubElement(spcAft_el, f'{{{nsmap}}}spcPts').set("val", str(space_after * 100))
            lnSpc = etree.SubElement(pPr, f'{{{nsmap}}}lnSpc')
            spcPct = etree.SubElement(lnSpc, f'{{{nsmap}}}spcPct')
            spcPct.set("val", "85000")
            if not text:
                return
            r_el = etree.SubElement(p_el, f'{{{nsmap}}}r')
            attrib = {"lang": "en-US", "dirty": "0"}
            if bold is not None:
                attrib["b"] = "1" if bold else "0"
            if size_pt:
                attrib["sz"] = str(int(size_pt * 100))
            rPr = etree.SubElement(r_el, f'{{{nsmap}}}rPr', attrib=attrib)
            if color_rgb:
                sf = etree.SubElement(rPr, f'{{{nsmap}}}solidFill')
                etree.SubElement(sf, f'{{{nsmap}}}srgbClr', attrib={"val": color_rgb})
            t_el = etree.SubElement(r_el, f'{{{nsmap}}}t')
            t_el.text = text

        def _clear_and_fill(txBody, lines):
            for p in txBody.findall(f'{{{nsmap}}}p'):
                txBody.remove(p)
            for item in lines:
                if len(item) == 4:
                    text, bold, size_pt, color_rgb = item
                    _add_para(txBody, text, bold, size_pt, color_rgb, 0, 0)
                elif len(item) == 6:
                    text, bold, size_pt, color_rgb, spc_before, spc_after = item
                    _add_para(txBody, text, bold, size_pt, color_rgb, spc_before, spc_after)
                else:
                    _add_para(txBody, item[0] if item else "", None, None, None, 0, 0)

        def fill_placeholder(ph_idx: int, lines: list):
            try:
                ph = slide.placeholders[ph_idx]
                _clear_and_fill(ph.text_frame._txBody, lines)
            except KeyError:
                logger.warning(f"Placeholder {ph_idx} not found")

        def fill_placeholder_resized(ph_idx: int, lines: list, size_pt: float, top: int, left: int, width: int, min_h: int, max_h: int = 0):
            try:
                ph = slide.placeholders[ph_idx]
                texts = [item[0] if isinstance(item, (list, tuple)) else "" for item in lines]
                est_h = _estimate_h(texts, size_pt, width)
                new_h = max(min_h, est_h)
                if max_h:
                    new_h = min(new_h, max_h)
                ph.top = top
                ph.left = left
                ph.width = width
                ph.height = new_h
                _apply_autofit(ph.text_frame._txBody)
                _clear_and_fill(ph.text_frame._txBody, lines)
                return new_h
            except KeyError:
                logger.warning(f"Placeholder {ph_idx} not found")
                return min_h

        def _get_font_size(item_count: int, base_min: int, base_max: int) -> int:
            if item_count <= 4:
                return base_max
            elif item_count <= 8:
                return base_max - 1
            elif item_count <= 12:
                return base_max - 2
            else:
                return base_min

        business_skills = candidate_data.get("business_skills", [])
        tech_skills = candidate_data.get("tech_skills", [])
        total_skill_items = len(business_skills) + len(tech_skills)
        skill_header_size = _get_font_size(total_skill_items, 8, 11)
        skill_item_size = _get_font_size(total_skill_items, 7, 10)

        languages = candidate_data.get("languages", [])
        industry = candidate_data.get("industry_experience", [])
        certs = candidate_data.get("certifications", [])
        education = candidate_data.get("education", [])
        total_right_items = len(languages) + len(industry) + len(certs) + len(education)
        right_header_size = _get_font_size(total_right_items, 8, 11)
        right_item_size = _get_font_size(total_right_items, 7, 10)

        summary_paras = candidate_data.get("summary_paras", [])
        total_summary_chars = sum(len(p) for p in summary_paras)
        if total_summary_chars < 300:
            summary_size = 11
        elif total_summary_chars < 500:
            summary_size = 10
        else:
            summary_size = 9

        exp_paras = [p for p in candidate_data.get("relevant_exp_paras", []) if p]
        total_exp_chars = sum(len(p) for p in exp_paras)
        if total_exp_chars < 400:
            exp_size = 10
        elif total_exp_chars < 700:
            exp_size = 9
        else:
            exp_size = 8

        truncated_summary = []
        for para in summary_paras:
            truncated = _truncate_to_fit(para, summary_size, PH26_WIDTH, 6)
            truncated_summary.append(truncated)

        truncated_exp = []
        for para in exp_paras:
            truncated = _truncate_to_fit(para, exp_size, EXP_TB_WIDTH, 5)
            truncated_exp.append(truncated)

        fill_placeholder(13, [
            (_sanitize(candidate_data.get("name_large", "")), None, 18, None),
        ])

        fill_placeholder(25, [
            (_sanitize(candidate_data.get("name_card", "")), True, 11, None),
            (_sanitize(candidate_data.get("role", "")), False, 9, None),
            (_sanitize(candidate_data.get("location", "")), False, 9, None),
        ])

        summary_lines = [
            (_sanitize(para), False, summary_size, None, 0, 0)
            for para in truncated_summary
        ]
        fill_placeholder_resized(
            26, summary_lines, summary_size,
            PH26_TOP, PH26_LEFT, PH26_WIDTH,
            PH26_MIN_H, PH26_MAX_H
        )

        skill_lines = [("Business Skills", True, skill_header_size, None, 0, 0)]
        for s in business_skills:
            skill_lines.append((_sanitize(s[:40]), False, skill_item_size, None, 0, 0))
        skill_lines.append(("", None, 4, None, 0, 0))
        skill_lines.append(("Technology Skills", True, skill_header_size, None, 0, 0))
        for item in tech_skills:
            if isinstance(item, (list, tuple)) and len(item) == 2:
                label, detail = _sanitize(item[0][:20]), _sanitize(item[1][:50])
            else:
                label, detail = _sanitize(item)[:20], ""
            skill_lines.append((f"{label}: {detail}", False, skill_item_size, None, 0, 0))

        skill_texts = [item[0] if isinstance(item, (list, tuple)) else "" for item in skill_lines]
        skill_est_h = _estimate_h(skill_texts, skill_item_size, PH16_WIDTH)
        skill_new_h = max(PH16_MIN_H, min(skill_est_h, LEFT_BOTTOM_LIMIT - PH16_TOP))
        try:
            ph16 = slide.placeholders[16]
            ph16.top = PH16_TOP
            ph16.left = PH16_LEFT
            ph16.width = PH16_WIDTH
            ph16.height = skill_new_h
            _apply_autofit(ph16.text_frame._txBody)
            _clear_and_fill(ph16.text_frame._txBody, skill_lines)
        except KeyError:
            pass

        right_lines = [("Languages", True, right_header_size, None, 0, 0)]
        for lang in languages:
            right_lines.append((_sanitize(lang[:40]), False, right_item_size, None, 0, 0))
        right_lines.append(("", None, 4, None, 0, 0))
        right_lines.append(("Industry Experience", True, right_header_size, None, 0, 0))
        for exp in industry:
            right_lines.append((_sanitize(exp[:40]), False, right_item_size, None, 0, 0))
        right_lines.append(("", None, 4, None, 0, 0))
        if certs:
            right_lines.append(("Certifications", True, right_header_size, None, 0, 0))
            for cert in certs:
                right_lines.append((_sanitize(cert[:50]), False, right_item_size, None, 0, 0))
            right_lines.append(("", None, 4, None, 0, 0))
        right_lines.append(("Education", True, right_header_size, None, 0, 0))
        for edu in education:
            right_lines.append((_sanitize(edu[:50]), False, right_item_size, None, 0, 0))

        right_texts = [item[0] if isinstance(item, (list, tuple)) else "" for item in right_lines]
        right_est_h = _estimate_h(right_texts, right_item_size, PH24_WIDTH)
        right_new_h = max(PH24_MIN_H, min(right_est_h, LEFT_BOTTOM_LIMIT - PH24_TOP))
        try:
            ph24 = slide.placeholders[24]
            ph24.top = PH24_TOP
            ph24.left = PH24_LEFT
            ph24.width = PH24_WIDTH
            ph24.height = right_new_h
            _apply_autofit(ph24.text_frame._txBody)
            _clear_and_fill(ph24.text_frame._txBody, right_lines)
        except KeyError:
            pass

        exp_lines = [
            (_sanitize(para), False, exp_size, None, 0, 2)
            for para in truncated_exp
        ]

        exp_textbox = None
        for shape in slide.shapes:
            if shape.name == "Text Placeholder 10" and shape.shape_type == 17:
                exp_textbox = shape
                break

        if exp_textbox is not None:
            text_list = [item[0] if isinstance(item, (list, tuple)) else "" for item in exp_lines]
            estimated_h = _estimate_h(text_list, exp_size, EXP_TB_WIDTH)
            new_exp_h = max(int(0.4 * 914400), min(estimated_h, EXP_TB_MAX_H))
            exp_textbox.top = EXP_TB_TOP
            exp_textbox.left = EXP_TB_LEFT
            exp_textbox.width = EXP_TB_WIDTH
            exp_textbox.height = new_exp_h
            _apply_autofit(exp_textbox.text_frame._txBody)
            _clear_and_fill(exp_textbox.text_frame._txBody, exp_lines)

            new_group25_top = EXP_TB_TOP + new_exp_h + SECTION_GAP
            new_clients_top = new_group25_top + GROUP25_HEIGHT + SECTION_GAP
            new_clients_bottom = new_clients_top + CLIENTS_HEIGHT

            if new_clients_bottom > SLIDE_HEIGHT:
                shift = new_clients_bottom - SLIDE_HEIGHT
                new_group25_top = max(EXP_TB_TOP + new_exp_h + SECTION_GAP, new_group25_top - shift)
                new_clients_top = new_group25_top + GROUP25_HEIGHT + SECTION_GAP

            for shape in slide.shapes:
                if shape.name == "Group 25" and shape.shape_type == 6:
                    shape.top = new_group25_top
                    shape.left = GROUP25_LEFT
                    shape.width = GROUP25_WIDTH
                    shape.height = GROUP25_HEIGHT
                    break

            try:
                clients_ph = slide.placeholders[28]
                clients_ph.top = new_clients_top
                clients_ph.left = CLIENTS_LEFT
                clients_ph.width = CLIENTS_WIDTH
                clients_ph.height = CLIENTS_HEIGHT
            except KeyError:
                pass

        clients_str = _sanitize(candidate_data.get("clients", ""))
        if len(clients_str) > 120:
            clients_str = clients_str[:117] + "..."
        clients_size = 8 if len(clients_str) < 80 else 7
        fill_placeholder(28, [
            (clients_str, False, clients_size, None, 0, 0),
        ])

        prs.save(out_path)

        with open(out_path, "rb") as f:
            return f.read()
    finally:
        try:
            os.unlink(out_path)
        except Exception:
            pass


async def _get_parsed_resume_for_candidate(session: AsyncSession, candidate: Candidate) -> Optional[ParsedResume]:
    from sqlalchemy import func
    if candidate.email and "@placeholder.com" not in candidate.email and "@noemail.vaspp.com" not in candidate.email:
        result = await session.execute(
            select(ParsedResume).where(ParsedResume.email == candidate.email).limit(1)
        )
        pr = result.scalars().first()
        if pr:
            return pr

    if candidate.first_name and candidate.last_name:
        result = await session.execute(
            select(ParsedResume).where(
                and_(
                    func.lower(ParsedResume.first_name) == candidate.first_name.lower(),
                    func.lower(ParsedResume.last_name) == candidate.last_name.lower(),
                )
            ).limit(1)
        )
        pr = result.scalars().first()
        if pr:
            return pr

    return None


class GenerateDeloitteResponse(BaseModel):
    cv_id: str
    deloitte_pptx_url: str
    message: str


@router.post("/{candidate_id}/cvs/{cv_id}/generate-deloitte", response_model=GenerateDeloitteResponse)
async def generate_deloitte_resume(
    candidate_id: UUID,
    cv_id: UUID,
    session: AsyncSession = Depends(get_db_session),
):
    candidate = await session.get(Candidate, candidate_id)
    if not candidate:
        raise HTTPException(status_code=404, detail="Candidate not found")

    cv_result = await session.execute(
        select(CandidateCV).where(
            and_(CandidateCV.id == cv_id, CandidateCV.candidate_id == candidate_id)
        )
    )
    cv = cv_result.scalar_one_or_none()
    if not cv:
        raise HTTPException(status_code=404, detail="CV not found")

    template_path = os.getenv(
        "DELOITTE_TEMPLATE_PATH",
        str(Path(__file__).resolve().parents[4] / "Deloitte_Template.pptx")
    )

    try:
        parsed = await _get_parsed_resume_for_candidate(session, candidate)

        cv_text = await _fetch_cv_text(cv.file_url)
        logger.info("cv_text_extracted", length=len(cv_text), candidate_id=str(candidate_id))

        candidate_data = None

        try:
            ai_ml_url = os.getenv("HR_APP_AI_SERVICE_URL", "http://ai-ml-service:8001")
            async with httpx.AsyncClient(timeout=60.0) as client:
                ai_response = await client.post(
                    f"{ai_ml_url}/api/v1/deloitte-parse",
                    json={
                        "cv_text": cv_text,
                        "candidate_id": str(candidate_id),
                        "candidate_location": str(candidate.location or ""),
                        "candidate_role": str(candidate.current_title or ""),
                    }
                )
                if ai_response.status_code == 200:
                    ai_data = ai_response.json()
                    if ai_data:
                        candidate_data = ai_data
                        logger.info("ai_ml_deloitte_parse_success", candidate_id=str(candidate_id))
        except Exception as e:
            logger.warning("ai_ml_deloitte_parse_failed", error=str(e), candidate_id=str(candidate_id))

        if candidate_data is None:
            logger.info("deloitte_regex_fallback", candidate_id=str(candidate_id))
            candidate_data = _parse_cv_text(cv_text, candidate, parsed)

        logger.info("candidate_data_built",
                    name=candidate_data.get("name_large"),
                    role=candidate_data.get("role"),
                    location=candidate_data.get("location"),
                    business_skills_count=len(candidate_data.get("business_skills", [])),
                    tech_skills_count=len(candidate_data.get("tech_skills", [])),
                    education_count=len(candidate_data.get("education", [])),
                    cert_count=len(candidate_data.get("certifications", [])),
                    clients=candidate_data.get("clients"),
                    candidate_id=str(candidate_id))

        pptx_bytes = _generate_pptx_bytes(candidate_data, template_path)

        name = candidate_data.get("name_large", f"{candidate.first_name}_{candidate.last_name}").replace(" ", "_")
        filename = f"{name}_Deloitte.pptx"

        if cv.deloitte_pptx_url:
            try:
                await delete_deloitte_pptx_from_storage(cv.deloitte_pptx_url)
            except Exception:
                pass

        pptx_url = await upload_deloitte_pptx(pptx_bytes, filename)
        if not pptx_url:
            raise HTTPException(status_code=500, detail="Failed to upload generated PPTX")

        cv.deloitte_pptx_url = pptx_url
        await session.commit()
        await session.refresh(cv)

        try:
            from src.core.redis import get_redis_pool
            redis = await get_redis_pool()
            await redis.delete(f"hr_app:candidate_profile:{candidate_id}")
        except Exception:
            pass

        return GenerateDeloitteResponse(
            cv_id=str(cv_id),
            deloitte_pptx_url=pptx_url,
            message="Deloitte resume generated successfully",
        )

    except HTTPException:
        raise
    except Exception as e:
        logger.error("deloitte_generate_failed", error=str(e), candidate_id=str(candidate_id), cv_id=str(cv_id))
        raise HTTPException(status_code=500, detail=f"Generation failed: {str(e)}")


@router.delete("/{candidate_id}/cvs/{cv_id}/deloitte", status_code=204)
async def delete_deloitte_resume(
    candidate_id: UUID,
    cv_id: UUID,
    session: AsyncSession = Depends(get_db_session),
):
    cv_result = await session.execute(
        select(CandidateCV).where(
            and_(CandidateCV.id == cv_id, CandidateCV.candidate_id == candidate_id)
        )
    )
    cv = cv_result.scalar_one_or_none()
    if not cv:
        raise HTTPException(status_code=404, detail="CV not found")

    if not cv.deloitte_pptx_url:
        raise HTTPException(status_code=404, detail="No Deloitte resume found for this CV")

    try:
        await delete_deloitte_pptx_from_storage(cv.deloitte_pptx_url)
    except Exception:
        pass

    cv.deloitte_pptx_url = None
    await session.commit()

    try:
        from src.core.redis import get_redis_pool
        redis = await get_redis_pool()
        await redis.delete(f"hr_app:candidate_profile:{candidate_id}")
    except Exception:
        pass