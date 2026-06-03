from pptx import Presentation
import lxml.etree as etree
import os

t = os.getenv('DELOITTE_TEMPLATE_PATH', '/app/Deloitte_Template.pptx')
prs = Presentation(t)
slide = prs.slides[0]

print('=== SLIDE MASTER XML - SLIDE NUMBER FIELDS ===')
xml_str = etree.tostring(slide.slide_layout.slide_master._element, pretty_print=True).decode()
for i, line in enumerate(xml_str.split('\n')):
    low = line.lower()
    if 'fld' in low or 'sldnum' in low or 'slidenum' in low:
        print(str(i) + ': ' + line)

print()
print('=== ALL MASTER SHAPES ===')
for shape in slide.slide_layout.slide_master.shapes:
    print('name=' + repr(shape.name) + ' type=' + str(shape.shape_type))
    if shape.has_text_frame:
        print('  text=' + repr(shape.text_frame.text.strip()))
        xml = etree.tostring(shape._element, pretty_print=True).decode()
        for line in xml.split('\n'):
            low = line.lower()
            if 'fld' in low or 'sldnum' in low:
                print('  XML: ' + line)
